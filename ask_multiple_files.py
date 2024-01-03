import streamlit as st
import tempfile
import numpy as np
import pandas as pd
import io
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.document_loaders.csv_loader import CSVLoader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import HuggingFaceInstructEmbeddings, OpenAIEmbeddings, HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
from langchain.chat_models import ChatOpenAI
from langchain.llms import HuggingFaceHub, OpenAI, CTransformers
from langchain.agents import create_csv_agent
import os
from pptx import Presentation

def load_csv_file(tmp_file_path):
    loader = CSVLoader(file_path=tmp_file_path, encoding="utf-8", csv_args={
            'delimiter': ','})
    data = loader.load()
    return data

def convert_excel_to_csv(excel_file):
    try:
        df = pd.read_excel(io.BytesIO(excel_file.read()))
        csv_file_path = excel_file.name.replace('.xlsx', '.csv')
        df.to_csv(csv_file_path, index=False)
        return csv_file_path
    
    except Exception as e:
        st.error(f"Error converting Excel to CSV: {e}")
        return None

def csv_agent_function(csv_file_path, user_question):
    agent = create_csv_agent(
        OpenAI(temperature=0), csv_file_path, verbose=True)
    if user_question is not None and user_question != "":
        st.write(agent.run(user_question))



def get_ppt_text(ppt_files):
    text = ""
    for ppt_file in ppt_files:
        presentation = Presentation(ppt_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text


def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=700,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks


def get_vectorstore_pdf_pptx(text_chunks, existing_vectorstore=None):
    embeddings = OpenAIEmbeddings()
    #embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
    if existing_vectorstore is not None:
        existing_vectorstore.save_local("faiss_index")
        existing_vectorstore = FAISS.load_local("faiss_index", embeddings)
    else:
        vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

#ToDo: add csv/excel to uploaded files
def save_uploaded_files(pdf_files, ppt_files):
    UPLOAD_FOLDER = "C:/Users/tschebesta/Desktop/MLA/ask-multiple-pdfs/uploaded_files"
    for pdf_file in pdf_files:
        save_path = os.path.join(UPLOAD_FOLDER, pdf_file.name)
        with open(save_path, "wb") as f:
            f.write(pdf_file.read())

    for ppt_file in ppt_files:
        save_path = os.path.join(UPLOAD_FOLDER, ppt_file.name)
        with open(save_path, "wb") as f:
            f.write(ppt_file.read())
        print(f"Saved {ppt_file.name} to {save_path}")



def get_conversation_chain_pdf_pptx(vectorstore):
    # llm = llm = CTransformers(model="llama-2-7b-chat.Q4_K_M.gguf", model_type='llama')
    llm = ChatOpenAI()
    #llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain


def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)


def main():
    load_dotenv()

    if os.getenv("OPENAI_API_KEY") is None or os.getenv("OPENAI_API_KEY") == "":
        print("OPENAI_API_KEY is not set")
        exit(1)

    st.set_page_config(page_title="Deloitte Private GPT",
                       page_icon=":milky_way:")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None
         
    st.session_state.vectorstore = None

    st.header("Deloitte Private GPT :milky_way:")
    user_question = st.text_input("Ask a question about your Files:")
    if user_question:
        handle_userinput(user_question)


    with st.sidebar:
        st.subheader("Your Files")
        pdf_docs = st.file_uploader("Upload PDF Files", type=['pdf'], accept_multiple_files=True)
        ppt_files = st.file_uploader("Upload PPTX Files", type=['pptx'], accept_multiple_files=True)
        csv_files = st.file_uploader("Upload CSV Files", type=['csv'], accept_multiple_files=True)
        excel_files = st.file_uploader("Upload Excel Files", type=['xlsx'], accept_multiple_files=True)
        if st.button("Process"):

            if pdf_docs or ppt_files:  
                with st.spinner("Processing Files"):
                    raw_text = ""

                    if pdf_docs:
                        raw_text += get_pdf_text(pdf_docs)
                    if ppt_files:
                        raw_text += get_ppt_text(ppt_files)
                    
                    existing_vectorstore = st.session_state.vectorstore
                    text_chunks = get_text_chunks(raw_text)

                    if existing_vectorstore is not None:
                        vectorstore = get_vectorstore_pdf_pptx(text_chunks, existing_vectorstore)
                    else:
                        vectorstore = get_vectorstore_pdf_pptx(text_chunks)
                    st.session_state.vectorstore = vectorstore
                    st.session_state.conversation = get_conversation_chain_pdf_pptx(vectorstore)
                    #save_uploaded_files(pdf_docs, ppt_files)

            elif csv_files or excel_files:
                with st.spinner("Processing CSV/Excel Files"):
                    if csv_files:
                        for uploaded_file in csv_files:
                            print('sowwy csv function still in progress...')

                    if excel_files:
                        for uploaded_file in excel_files:
                            csv_file_path = convert_excel_to_csv(uploaded_file)
                            csv_agent_function(csv_file_path, user_question)
                #save_uploaded_files(csv_files)
            else:
                st.warning("Please upload files before clicking 'Process'")






if __name__ == '__main__':
    main()
